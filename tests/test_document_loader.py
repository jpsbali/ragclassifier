from io import BytesIO
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation

import src.document_loader as document_loader


def test_extract_text_from_txt() -> None:
    text = document_loader.extract_text_from_upload("sample.txt", b"Hello world")
    assert text == "Hello world"


def test_extract_text_from_docx_bytes() -> None:
    doc = Document()
    doc.add_paragraph("First line")
    doc.add_paragraph("Second line")
    buffer = BytesIO()
    doc.save(buffer)

    text = document_loader.extract_text_from_upload("sample.docx", buffer.getvalue())
    assert text == "First line\nSecond line"


def test_extract_text_from_pdf_uses_pypdf_reader(monkeypatch) -> None:
    class FakePage:
        def __init__(self, value: str) -> None:
            self._value = value

        def extract_text(self) -> str:
            return self._value

    class FakeReader:
        def __init__(self, file_obj: BytesIO) -> None:
            assert hasattr(file_obj, "read")
            self.pages = [FakePage("Page one"), FakePage("Page two")]

    monkeypatch.setattr(document_loader, "PdfReader", FakeReader)

    text = document_loader.extract_text_from_upload("sample.pdf", b"%PDF-FAKE")
    assert text == "Page one\nPage two"


def test_extract_text_for_real_rules_docx() -> None:
    project_root = Path(__file__).resolve().parents[1]
    doc_path = project_root / "ClassifyingRules.docx"
    text = document_loader.extract_text_from_upload(doc_path.name, doc_path.read_bytes())

    assert "RESTRICTED" in text
    assert "CONFIDENTIAL" in text
    assert "PUBLIC" in text


def test_extract_text_from_pptx_bytes() -> None:
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "Slide Title"

    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "First bullet point"
    p = tf.add_paragraph()
    p.text = "Second bullet point"

    buffer = BytesIO()
    prs.save(buffer)

    text = document_loader.extract_text_from_upload("sample.pptx", buffer.getvalue())
    assert text == "Slide Title\nFirst bullet point\nSecond bullet point"


def test_extract_text_from_xlsx_bytes() -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws["A1"] = "Header 1"
    ws["B1"] = "Header 2"
    ws.append(["Data 1A", "Data 1B"])
    ws.append(["Data 2A", None, "Data 2C"])  # Add a gap

    ws2 = wb.create_sheet("Sheet2")
    ws2.append(["Sheet 2 Data"])

    buffer = BytesIO()
    wb.save(buffer)

    text = document_loader.extract_text_from_upload("sample.xlsx", buffer.getvalue())
    assert text == "Header 1 | Header 2\nData 1A | Data 1B\nData 2A | Data 2C\nSheet 2 Data"
