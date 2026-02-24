from io import BytesIO

from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation


def extract_text_from_upload(filename: str, content_bytes: bytes) -> str:
    lower_name = filename.lower()
    if lower_name.endswith(".txt") or lower_name.endswith(".md"):
        return content_bytes.decode("utf-8", errors="ignore").strip()
    if lower_name.endswith(".pdf"):
        reader = PdfReader(BytesIO(content_bytes))
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n".join(pages).strip()
    if lower_name.endswith(".docx"):
        doc = Document(BytesIO(content_bytes))
        return "\n".join(p.text for p in doc.paragraphs).strip()
    if lower_name.endswith(".pptx"):
        prs = Presentation(BytesIO(content_bytes))
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        return "\n".join(text_runs).strip()
    if lower_name.endswith(".xlsx"):
        workbook = load_workbook(filename=BytesIO(content_bytes), read_only=True)
        text_parts = []
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows():
                row_text = [str(cell.value) for cell in row if cell.value is not None]
                if row_text:
                    text_parts.append(" | ".join(row_text))
        return "\n".join(text_parts).strip()

    # Fallback for unknown file types.
    return content_bytes.decode("utf-8", errors="ignore").strip()
